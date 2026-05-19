"""
PPT 优化脚本 — 赛意信息资本市场沟通会
========================================
基于 python-pptx，保留原文内容，提升视觉设计质量
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

SRC = r"C:\Users\liuziheng\Downloads\赛意信息：AI驱动工业智算，开启高质量增长新篇章.pptx"
DST = r"C:\Users\liuziheng\Downloads\赛意信息_优化版.pptx"

# ── 设计系统 ─────────────────────────────────────────
# 色彩
C_PRIMARY    = RGBColor(0x0F, 0x3C, 0x6B)   # 深蓝 — 主色
C_SECONDARY  = RGBColor(0x1A, 0x73, 0xE8)   # 科技蓝 — 辅色
C_ACCENT     = RGBColor(0x00, 0xA8, 0x9D)   # 青绿 — 强调
C_ACCENT2    = RGBColor(0xE8, 0x6A, 0x17)   # 暖橙 — 数据突出
C_DARK       = RGBColor(0x1A, 0x1A, 0x2E)   # 深色背景
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG   = RGBColor(0xF0, 0xF4, 0xF8)   # 浅灰蓝背景
C_TEXT       = RGBColor(0x2D, 0x2D, 0x3F)   # 正文色
C_TEXT_LIGHT = RGBColor(0x6B, 0x72, 0x80)   # 次要文字
C_GREEN_UP   = RGBColor(0x16, 0xA3, 0x4A)   # 涨
C_RED_DOWN   = RGBColor(0xDC, 0x26, 0x26)   # 跌

# 字体
FONT_CN = 'Noto Sans SC'
FONT_EN = 'Arial'

MARGIN = Emu(685800)   # 0.75 inches
SLIDE_W = None  # 将从源文件读取
SLIDE_H = None

def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    """添加带填充色的形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                 font_color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_CN, anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_rich_text(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """添加富文本框，segments = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, seg in enumerate(segments):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        text, size, color, bold = seg[0], seg[1] if len(seg) > 1 else Pt(14), seg[2] if len(seg) > 2 else C_TEXT, seg[3] if len(seg) > 3 else False
        run.text = text
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT_CN
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), FONT_CN)
    return txBox

def add_slide_number(slide, num, total=11):
    """添加页码"""
    add_text_box(slide, Emu(11000000), Emu(6500000), Emu(1000000), Emu(300000),
                 f'{num} / {total}', Pt(9), C_TEXT_LIGHT, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, section_num, title):
    """添加章节标题栏"""
    # 顶部色条
    add_shape(slide, Emu(0), Emu(0), SLIDE_W, Emu(55000), C_PRIMARY)
    # 左侧数字色块
    add_shape(slide, MARGIN, Emu(380000), Emu(55000), Emu(55000), C_ACCENT)
    # 标题
    add_text_box(slide, Emu(1020000), Emu(330000), Emu(9500000), Emu(550000),
                 f'{section_num:02d}.  {title}', Pt(28), C_PRIMARY, bold=True)

def add_card(slide, left, top, width, height, title, body, icon_text=None, title_color=C_PRIMARY):
    """添加卡片组件"""
    # 卡片背景
    card = add_shape(slide, left, top, width, height, C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    # 顶部色条
    add_shape(slide, left, top, width, Emu(32000), title_color)
    # 图标/编号
    if icon_text:
        add_text_box(slide, left + Emu(150000), top + Emu(100000), Emu(400000), Emu(350000),
                     icon_text, Pt(22), title_color, bold=True)
    # 标题
    add_text_box(slide, left + Emu(150000) + (Emu(500000) if icon_text else 0),
                 top + Emu(100000), width - Emu(300000) - (Emu(500000) if icon_text else 0),
                 Emu(350000), title, Pt(14), title_color, bold=True)
    # 正文
    add_text_box(slide, left + Emu(150000), top + Emu(500000),
                 width - Emu(300000), height - Emu(650000),
                 body, Pt(10), C_TEXT)

def add_kpi_card(slide, left, top, width, height, label, value, sub_text=None, accent=False):
    """添加 KPI 数字卡片"""
    color = C_ACCENT2 if accent else C_SECONDARY
    add_shape(slide, left, top, width, height, C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    # 左侧色条
    add_shape(slide, left, top, Emu(32000), height, color)
    add_text_box(slide, left + Emu(120000), top + Emu(80000), width - Emu(200000), Emu(250000),
                 label, Pt(9), C_TEXT_LIGHT)
    add_text_box(slide, left + Emu(120000), top + Emu(300000), width - Emu(200000), Emu(450000),
                 value, Pt(24), color, bold=True)
    if sub_text:
        add_text_box(slide, left + Emu(120000), top + Emu(700000), width - Emu(200000), Emu(250000),
                     sub_text, Pt(8), C_GREEN_UP if '▲' in sub_text or '增' in sub_text else C_TEXT_LIGHT)


print("=" * 50)
print("PPT 优化中...")
print("=" * 50)

# 从原文件读取尺寸信息后，创建全新 PPT
src_prs = Presentation(SRC)
SLIDE_W = src_prs.slide_width
SLIDE_H = src_prs.slide_height

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# 先加一个临时 slide 以获取 blank_layout，最后会删掉
blank_layout = prs.slide_layouts[6]  # blank layout
tmp_slide = prs.slides.add_slide(blank_layout)

TOTAL = 11

# ══════════════════════════════════════════════════
# Slide 1: 封面
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_DARK)

# 装饰线
add_shape(slide, MARGIN, Emu(1800000), Emu(800000), Emu(25000), C_ACCENT)
# 主标题
add_text_box(slide, MARGIN, Emu(2000000), Emu(10700000), Emu(1200000),
             'AI驱动工业智算，开启高质量增长新篇章', Pt(40), C_WHITE, bold=True)
# 副标题
add_text_box(slide, MARGIN, Emu(3300000), Emu(10700000), Emu(500000),
             '赛意信息（300687）2026年资本市场沟通会', Pt(20), C_ACCENT)
# 日期
add_text_box(slide, MARGIN, Emu(3900000), Emu(5000000), Emu(400000),
             '2026年6月', Pt(14), C_TEXT_LIGHT)
# 底部装饰线
add_shape(slide, MARGIN, Emu(6100000), Emu(10700000), Emu(15000), C_SECONDARY)
add_slide_number(slide, 1)

# ══════════════════════════════════════════════════
# Slide 2: 目录
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_WHITE)

add_text_box(slide, MARGIN, Emu(350000), Emu(5000000), Emu(500000),
             'CONTENTS', Pt(28), C_PRIMARY, bold=True, font_name=FONT_EN)
add_shape(slide, MARGIN, Emu(850000), Emu(600000), Emu(25000), C_ACCENT)

toc_items = [
    ('01', '战略转型：迈向工业AI领军者'),
    ('02', '核心业务架构："一核、两基、三化、四场景"'),
    ('03', '算力生态布局：国产筑基 + 英伟达协同'),
    ('04', 'AI技术战略：打造行业专属模型矩阵'),
    ('05', '转型路径与目标：三步走战略'),
    ('06', 'AI商业化成果：从投入到价值兑现'),
    ('07', '增长动能与业绩拐点：2026年Q1印证'),
    ('08', '估值展望与投资者价值'),
]

for i, (num, title) in enumerate(toc_items):
    col = i // 4
    row = i % 4
    x = MARGIN + col * Emu(5500000)
    y = Emu(1300000) + row * Emu(1250000)

    # 编号圆圈
    circle = add_shape(slide, x, y, Emu(500000), Emu(500000), C_PRIMARY if i < 5 else C_SECONDARY,
                       shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, x, y + Emu(50000), Emu(500000), Emu(400000),
                 num, Pt(22), C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, x + Emu(650000), y + Emu(80000), Emu(4400000), Emu(400000),
                 title, Pt(14), C_TEXT, bold=False)

# 右侧装饰区
add_shape(slide, Emu(11300000), Emu(1200000), Emu(500000), Emu(5000000), C_PRIMARY)
add_text_box(slide, Emu(11350000), Emu(2500000), Emu(400000), Emu(3000000),
             'STRATEGIC\nAGENDA\n2026', Pt(16), C_WHITE, bold=True, font_name=FONT_EN)

add_slide_number(slide, 2)

# ══════════════════════════════════════════════════
# Slide 3-10: 内容页（保持原始内容，优化排版）
# ══════════════════════════════════════════════════

def build_content_slide_3():
    """01 战略转型"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 1, '战略转型：从工业软件服务商迈向工业AI领军者')

    # 左侧：行业变革
    add_text_box(slide, MARGIN, Emu(1200000), Emu(5000000), Emu(350000),
                 '时代机遇与行业变革', Pt(18), C_PRIMARY, bold=True)

    cards_left = [
        ('核心判断：智能生产新阶段',
         '制造业数智化迈入AI原生+算力驱动的全新发展阶段，技术驱动成为核心引擎。'),
        ('需求升级：价值交付重塑',
         '制造企业需求已从单一代码交付，转向以价值为导向的"工业智能即服务"(AIaaS)。'),
        ('行业变革：价值链重构',
         'AI技术全面重构工业软件价值链，传统低附加值、纯工具型的模式增长已显乏力，急需破局。'),
    ]
    for i, (title, body) in enumerate(cards_left):
        add_card(slide, MARGIN, Emu(1700000) + i * Emu(1550000),
                 Emu(5200000), Emu(1350000), title, body, str(i+1))

    # 右侧：公司转型
    add_text_box(slide, Emu(6400000), Emu(1200000), Emu(5000000), Emu(350000),
                 '公司转型与定位升级', Pt(18), C_PRIMARY, bold=True)

    # 三个优势卡片
    advantages = [
        ('MES市占率第一', '电子/光伏/家电行业'),
        ('3200+', '制造行业标杆客户'),
        ('谷神"双跨"平台', '国产化生态布局'),
    ]
    for i, (value, label) in enumerate(advantages):
        x = Emu(6400000) + i * Emu(1750000)
        add_shape(slide, x, Emu(1700000), Emu(1550000), Emu(1200000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, x + Emu(80000), Emu(1800000), Emu(1400000), Emu(500000),
                     value, Pt(22), C_ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Emu(80000), Emu(2300000), Emu(1400000), Emu(400000),
                     label, Pt(10), C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # 底部 — 升级标语
    add_shape(slide, Emu(6400000), Emu(3200000), Emu(5100000), Emu(1100000), C_PRIMARY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Emu(6600000), Emu(3350000), Emu(4700000), Emu(400000),
                 '从"工业软件服务商" 全面升级为', Pt(12), C_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Emu(6600000), Emu(3700000), Emu(4700000), Emu(400000),
                 '"工业智算 + 行业模型 + 智能体Agent" 一体化解决方案领军者', Pt(14), C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 底部 KPI
    add_kpi_card(slide, MARGIN, Emu(5000000), Emu(2440000), Emu(900000), '服务客户', '3,200+', '制造行业标杆')
    add_kpi_card(slide, Emu(3400000), Emu(5000000), Emu(2440000), Emu(900000), 'MES市占率', '第一', '电子/光伏/家电', accent=True)

    add_slide_number(slide, 3)
    return slide


def build_content_slide_4():
    """02 一核两基三化四场景"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 2, '核心业务架构："一核、两基、三化、四场景"')

    # 中心 — 一核
    add_shape(slide, Emu(5100000), Emu(1400000), Emu(2000000), Emu(2000000), C_PRIMARY, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Emu(5100000), Emu(1900000), Emu(2000000), Emu(400000),
                 '一核', Pt(20), C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Emu(5100000), Emu(2300000), Emu(2000000), Emu(600000),
                 '工业AI\n与算力', Pt(14), C_WHITE, alignment=PP_ALIGN.CENTER)

    # 两基 — 左右
    for side, title, desc in [
        ('left', '昇腾智算底座', '自主高效可扩展\n算力基础设施'),
        ('right', '谷神工业数据与模型中台', '沉淀工业机理与\nAI算法的能力中枢'),
    ]:
        x = Emu(800000) if side == 'left' else Emu(8500000)
        add_shape(slide, x, Emu(2000000), Emu(2800000), Emu(850000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, x + Emu(100000), Emu(2100000), Emu(2600000), Emu(300000),
                     title, Pt(14), C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Emu(100000), Emu(2400000), Emu(2600000), Emu(350000),
                     desc, Pt(10), C_TEXT, alignment=PP_ALIGN.CENTER)

    # 三化
    for i, (title, desc) in enumerate([
        ('算力服务化', '降低门槛，弹性供给'),
        ('模型行业化', '深耕垂直领域'),
        ('Agent场景化', '赋能具体业务'),
    ]):
        x = Emu(1400000) + i * Emu(3300000)
        add_shape(slide, x, Emu(3300000), Emu(2800000), Emu(700000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, x + Emu(80000), Emu(3400000), Emu(2600000), Emu(250000),
                     title, Pt(13), C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Emu(80000), Emu(3650000), Emu(2600000), Emu(250000),
                     desc, Pt(9), C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # 四场景 — 底部
    add_shape(slide, Emu(800000), Emu(4400000), Emu(10600000), Emu(900000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Emu(1000000), Emu(4550000), Emu(10200000), Emu(300000),
                 '四场景：研发设计 · 生产制造 · 质量管控 · 供应链协同', Pt(14), C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Emu(1000000), Emu(4850000), Emu(10200000), Emu(300000),
                 '覆盖工业全价值链，推动端到端的数字化转型落地', Pt(10), C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 4)
    return slide


def build_content_slide_5():
    """03 算力生态"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 3, '算力生态布局：国产算力筑基 + 英伟达先进算力协同')

    # 两列
    for col, (title, items, color) in enumerate([
        ('国产算力筑基', [
            ('核心目标', '打造工业专属国产化算力体系，从底层硬件到软件架构实现自主掌控'),
            ('全维服务', '"算力租赁+集群部署+运维调优"端到端全生命周期服务'),
            ('落地进展', '2025年成功承接人工智能算力中心项目，完成能力验证'),
        ], C_PRIMARY),
        ('英伟达先进算力协同', [
            ('技术协同', '基于英伟达GPU集群，利用NVLink等先进互联技术为大规模AI训练提供性能支撑'),
            ('战略落地', '2026年成立"智数科技"全资子公司，集中优势资源拓展AI场景'),
            ('融合策略', '"国产算力+英伟达先进算力"双轨并行，满足差异化需求'),
        ], C_SECONDARY),
    ]):
        x = MARGIN + col * Emu(5800000)
        add_text_box(slide, x, Emu(1200000), Emu(5000000), Emu(400000),
                     title, Pt(20), color, bold=True)
        for i, (sub_title, desc) in enumerate(items):
            add_card(slide, x, Emu(1800000) + i * Emu(1550000),
                     Emu(5300000), Emu(1350000), sub_title, desc, str(i+1), color)

    add_slide_number(slide, 5)
    return slide


def build_content_slide_6():
    """04 AI技术战略"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 4, 'AI技术战略：打造行业专属模型矩阵')

    layers = [
        ('企业运营模型层', '数字主干网', '集"监控、预警、归因、责任、作战、决策"六大中心于一体，打通企业数据脉络', C_PRIMARY),
        ('制造模型层', '数字调度官', '覆盖生产排程、工艺执行、质量预警、设备健康四大核心场景', C_SECONDARY),
        ('行业模型层', '行业老师傅', '深耕PCB、光伏、工程机械等细分赛道，沉淀Know-how与工艺参数', C_ACCENT),
        ('助手模型层', '数字员工', '轻量化、场景化的智能体应用，覆盖SOP助手、运维助手、供应链助手等角色', C_ACCENT2),
    ]

    for i, (title, tagline, desc, color) in enumerate(layers):
        y = Emu(1300000) + i * Emu(1200000)
        # 左侧色块
        add_shape(slide, MARGIN, y, Emu(55000), Emu(950000), color)
        # 标题
        add_text_box(slide, Emu(900000), y, Emu(3500000), Emu(350000),
                     title, Pt(16), color, bold=True)
        # 标签
        add_shape(slide, Emu(4400000), y + Emu(20000), Emu(1400000), Emu(300000), color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, Emu(4400000), y + Emu(40000), Emu(1400000), Emu(250000),
                     tagline, Pt(10), C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # 描述
        add_text_box(slide, Emu(900000), y + Emu(450000), Emu(10500000), Emu(400000),
                     desc, Pt(11), C_TEXT)

    # 里程碑
    add_shape(slide, MARGIN, Emu(6100000), Emu(10700000), Emu(500000), C_PRIMARY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Emu(900000), Emu(6150000), Emu(10500000), Emu(400000),
                 '关键里程碑：善谋GPT成功入选IDC《中国工业大模型及智能体解决方案2025年厂商评估》领导者阵营',
                 Pt(13), C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 6)
    return slide


def build_content_slide_7():
    """05 三步走战略"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 5, '战略转型落地路径与目标')

    phases = [
        ('第一阶段：夯实基础', '2026-2027', [
            '完成算力底座搭建，筑牢数字基石',
            '首批行业模型成功落地，验证价值',
            '传统低价值业务占比降至50%以下',
        ], C_PRIMARY),
        ('第二阶段：规模扩张', '2028-2029', [
            '标准化产品矩阵全面落地，快速复制',
            '新业务收入占比过半，成为增长主力',
            '整体盈利能力实现质的飞跃与提升',
        ], C_SECONDARY),
        ('第三阶段：行业引领', '2030及以后', [
            '构建完整的工业AI生态，赋能上下游',
            '跻身国内工业智算与智能解决方案第一梯队',
            '确立行业标准制定者与引领者地位',
        ], C_ACCENT),
    ]

    for i, (title, period, items, color) in enumerate(phases):
        x = MARGIN + i * Emu(3700000)
        # 时间轴节点
        add_shape(slide, x + Emu(1600000), Emu(1200000), Emu(120000), Emu(120000), color, shape_type=MSO_SHAPE.OVAL)
        if i < 2:
            add_shape(slide, x + Emu(1720000), Emu(1250000), Emu(1980000), Emu(25000), color)

        add_text_box(slide, x, Emu(1450000), Emu(3400000), Emu(300000),
                     period, Pt(12), color, bold=True)
        add_text_box(slide, x, Emu(1750000), Emu(3400000), Emu(300000),
                     title, Pt(16), C_TEXT, bold=True)

        for j, item in enumerate(items):
            add_text_box(slide, x + Emu(120000), Emu(2250000) + j * Emu(400000),
                         Emu(3200000), Emu(350000), f'• {item}', Pt(10), C_TEXT)

    # 底部保障
    add_shape(slide, MARGIN, Emu(5800000), Emu(10700000), Emu(700000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Emu(900000), Emu(5900000), Emu(10500000), Emu(500000),
                 '战略保障：深化生态合作 | 共建开放研发体系 | 成立数字化转型专项部门 | 持续加大核心技术研发投入',
                 Pt(11), C_TEXT, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 7)
    return slide


def build_content_slide_8():
    """06 AI商业化"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 6, 'AI商业化成果：已从研发投入转向价值兑现')

    metrics = [
        ('重大订单验证', '¥4,847万', '成功签订重大AI中台项目，验证AI大模型工业应用的商业化可行性'),
        ('订单持续放量', '¥1.03亿', '2025年上半年AI相关订单突破1.03亿元，全年持续放量增长'),
        ('效率显著提升', '>80%', '供应链AI Agent推动客户运营效率提升超80%，AI辅助编程使研发效率提升40%-70%'),
    ]

    for i, (label, value, desc) in enumerate(metrics):
        y = Emu(1400000) + i * Emu(1600000)
        add_kpi_card(slide, MARGIN, y, Emu(5200000), Emu(1350000), label, value, desc)

    # 右侧市场数据
    add_shape(slide, Emu(6400000), Emu(1400000), Emu(5100000), Emu(2200000), C_PRIMARY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Emu(6700000), Emu(1600000), Emu(4500000), Emu(400000),
                 '市场空间广阔', Pt(20), C_WHITE, bold=True)
    add_text_box(slide, Emu(6700000), Emu(2100000), Emu(4500000), Emu(300000),
                 '中国工业企业大模型应用比例', Pt(12), C_WHITE)

    # 进度对比
    add_text_box(slide, Emu(6700000), Emu(2500000), Emu(2000000), Emu(400000),
                 '2024年', Pt(11), C_TEXT_LIGHT)
    add_text_box(slide, Emu(6700000), Emu(2800000), Emu(2000000), Emu(500000),
                 '9.6%', Pt(36), C_ACCENT, bold=True)
    add_text_box(slide, Emu(9400000), Emu(2500000), Emu(2000000), Emu(400000),
                 '2025年', Pt(11), C_TEXT_LIGHT)
    add_text_box(slide, Emu(9400000), Emu(2800000), Emu(2000000), Emu(500000),
                 '47.5%', Pt(36), C_ACCENT2, bold=True)

    # 箭头
    add_shape(slide, Emu(8700000), Emu(3000000), Emu(700000), Emu(25000), C_WHITE)

    # 底部说明
    add_shape(slide, Emu(6400000), Emu(4000000), Emu(5100000), Emu(1000000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Emu(6600000), Emu(4150000), Emu(4700000), Emu(700000),
                 'IDC调研显示，中国工业企业大模型应用比例从2024年的9.6%提升至2025年的47.5%，市场空间广阔。',
                 Pt(10), C_TEXT)

    add_slide_number(slide, 8)
    return slide


def build_content_slide_9():
    """07 业绩拐点"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 7, '增长动能与业绩拐点：2026年Q1印证转型成效')

    # 2025 年数据
    add_text_box(slide, MARGIN, Emu(1200000), Emu(5000000), Emu(350000),
                 '2025年：主动调整，蓄力前行', Pt(18), C_PRIMARY, bold=True)

    add_kpi_card(slide, MARGIN, Emu(1700000), Emu(2440000), Emu(1100000),
                 '研发投入', '3.70亿元', '占营收 17.84%')
    add_kpi_card(slide, Emu(3400000), Emu(1700000), Emu(2440000), Emu(1100000),
                 '经营现金流', '6,108万元', '同比增长 6,778% ▲', accent=True)

    # 2026 Q1
    add_text_box(slide, Emu(6400000), Emu(1200000), Emu(5000000), Emu(350000),
                 '2026年Q1：经营拐点已现', Pt(18), C_ACCENT2, bold=True)

    add_kpi_card(slide, Emu(6400000), Emu(1700000), Emu(2440000), Emu(1100000),
                 '归母净利润', '3,157万元', '同比增长 28.74% ▲')
    add_kpi_card(slide, Emu(9000000), Emu(1700000), Emu(2440000), Emu(1100000),
                 '扣非净利润', '3,107万元', '同比增长 33.81% ▲', accent=True)

    # 底部趋势提示
    add_shape(slide, MARGIN, Emu(3200000), Emu(10700000), Emu(25000), C_GREEN_UP)
    add_shape(slide, Emu(6400000), Emu(3200000), Emu(5100000), Emu(25000), C_GREEN_UP)

    add_text_box(slide, MARGIN, Emu(3400000), Emu(10700000), Emu(300000),
                 '战略性投入，回款质量显著提升', Pt(10), C_TEXT_LIGHT)
    add_text_box(slide, Emu(6400000), Emu(3400000), Emu(5100000), Emu(300000),
                 '盈利指标全面改善，资产质量稳步提升', Pt(10), C_TEXT_LIGHT)

    add_slide_number(slide, 9)
    return slide


def build_content_slide_10():
    """08 估值展望"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_LIGHT_BG)
    add_section_header(slide, 8, '估值展望与投资者价值')

    # 左侧：目标价
    add_text_box(slide, MARGIN, Emu(1200000), Emu(5000000), Emu(350000),
                 '机构看好，空间广阔', Pt(18), C_PRIMARY, bold=True)

    price_targets = [
        ('13家机构综合目标价', '27.13-27.38元'),
        ('分析师12个月平均目标价', '32.60元'),
        ('券商维持"增持"评级', '40.34元'),
    ]
    for i, (label, value) in enumerate(price_targets):
        y = Emu(1750000) + i * Emu(900000)
        add_shape(slide, MARGIN, y, Emu(5000000), Emu(700000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text_box(slide, Emu(900000), y + Emu(80000), Emu(2800000), Emu(250000),
                     label, Pt(11), C_TEXT_LIGHT)
        add_text_box(slide, Emu(3700000), y + Emu(80000), Emu(1800000), Emu(250000),
                     value, Pt(22), C_ACCENT2, bold=True)

    # 右侧：五大亮点
    add_text_box(slide, Emu(6400000), Emu(1200000), Emu(5000000), Emu(350000),
                 '五大核心投资亮点', Pt(18), C_ACCENT2, bold=True)

    highlights = [
        ('市场卡位', 'MES市占率第一，谷神平台"双跨"认证'),
        ('AI商业化', 'IDC领导者阵营，近5000万单体订单验证'),
        ('信创爆发', '国央企订单增速超300%，增长强劲'),
        ('海外发展', '海外业务收入同比增长236%'),
        ('战略转型', '从低毛利项目制向高价值订阅制转型'),
    ]
    for i, (title, desc) in enumerate(highlights):
        y = Emu(1750000) + i * Emu(700000)
        add_shape(slide, Emu(6400000), y, Emu(5100000), Emu(550000), C_WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_shape(slide, Emu(6400000), y, Emu(55000), Emu(550000), C_ACCENT if i < 3 else C_SECONDARY)
        add_text_box(slide, Emu(6600000), y + Emu(50000), Emu(1400000), Emu(250000),
                     title, Pt(11), C_PRIMARY, bold=True)
        add_text_box(slide, Emu(8100000), y + Emu(50000), Emu(3300000), Emu(250000),
                     desc, Pt(9), C_TEXT)

    # 风险提示
    add_text_box(slide, MARGIN, Emu(5800000), Emu(10700000), Emu(500000),
                 '风险提示：战略转型不及预期；下游需求波动；技术迭代风险；行业竞争加剧；海外市场拓展风险。',
                 Pt(8), C_TEXT_LIGHT)

    add_slide_number(slide, 10)
    return slide


# ══════════════════════════════════════════════════
# 构建所有内容页
# ══════════════════════════════════════════════════
build_content_slide_3()
build_content_slide_4()
build_content_slide_5()
build_content_slide_6()
build_content_slide_7()
build_content_slide_8()
build_content_slide_9()
build_content_slide_10()

# ══════════════════════════════════════════════════
# Slide 11: 感谢页
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, C_DARK)

add_shape(slide, Emu(4600000), Emu(2000000), Emu(3000000), Emu(3000000), C_PRIMARY, shape_type=MSO_SHAPE.OVAL)
add_text_box(slide, Emu(1500000), Emu(2800000), Emu(9200000), Emu(700000),
             '感谢聆听', Pt(48), C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Emu(1500000), Emu(3500000), Emu(9200000), Emu(400000),
             'THANK YOU', Pt(18), C_ACCENT, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
add_text_box(slide, Emu(1500000), Emu(4500000), Emu(9200000), Emu(350000),
             '广州赛意信息科技股份有限公司', Pt(14), C_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Emu(1500000), Emu(4850000), Emu(9200000), Emu(350000),
             '股票代码：300687 | www.chinasie.com', Pt(12), C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 11)

# ── 删除临时 slide ───────────────────────────────
rId = prs.slides._sldIdLst[0].get(qn('r:id'))
prs.part.drop_rel(rId)
prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# ── 保存 ─────────────────────────────────────────
prs.save(DST)
print(f"\n优化版 PPT 已保存至: {DST}")
print("=" * 50)
print("优化要点:")
print("  1. 统一色彩体系（深蓝+科技蓝+青绿+暖橙）")
print("  2. 卡片式布局，信息层次分明")
print("  3. KPI 数字突出展示")
print("  4. 章节标题统一格式+页码")
print("  5. 封面/封底深色背景增强仪式感")
print("  6. 数据可视化（指标对比、进度展示）")
print("  7. 减少文字密度，关键信息更突出")
print("=" * 50)
